#!/usr/bin/env python3
"""Build the editable VLA-Corrector presentation deck and web previews.

The PPTX is intended for manual editing. The GitHub Pages presentation viewer
uses exported slide PNGs generated from the same source content.
"""

from __future__ import annotations

import subprocess
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
DOCS = ROOT / "docs"
ASSETS = DOCS / "assets"
PRESENTATION = ASSETS / "presentation"
SLIDES = PRESENTATION / "slides"
THUMBS = PRESENTATION / "thumbnails"

W, H = 1600, 900
BLUE = "#2f5f9f"
BLUE_DARK = "#173b68"
BLUE_SOFT = "#edf3fb"
TEXT = "#17202a"
MUTED = "#5f6b7a"
LINE = "#dbe3ee"
PAPER = "#fbfcfe"
GREEN = "#2f8f5b"
RED = "#b94a48"

FONT_REG = "/usr/share/fonts/truetype/msttcorefonts/Arial.ttf"
FONT_BOLD = "/usr/share/fonts/truetype/msttcorefonts/arialbd.ttf"
FONT_GEORGIA = "/usr/share/fonts/truetype/msttcorefonts/georgiab.ttf"


def rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.lstrip("#")
    return RGBColor(*(int(hex_color[i : i + 2], 16) for i in (0, 2, 4)))


def pil_font(size: int, bold: bool = False, serif: bool = False) -> ImageFont.FreeTypeFont:
    path = FONT_GEORGIA if serif else (FONT_BOLD if bold else FONT_REG)
    return ImageFont.truetype(path, size=size)


def hex_rgb(hex_color: str) -> tuple[int, int, int]:
    hex_color = hex_color.lstrip("#")
    return tuple(int(hex_color[i : i + 2], 16) for i in (0, 2, 4))


def text_wrap(draw: ImageDraw.ImageDraw, text: str, font: ImageFont.FreeTypeFont, max_width: int) -> list[str]:
    lines: list[str] = []
    for paragraph in text.split("\n"):
        words = paragraph.split()
        line = ""
        for word in words:
            candidate = f"{line} {word}".strip()
            if draw.textbbox((0, 0), candidate, font=font)[2] <= max_width:
                line = candidate
            else:
                if line:
                    lines.append(line)
                line = word
        if line:
            lines.append(line)
    return lines


def draw_wrapped(draw, xy, text, font, fill, max_width, line_gap=8):
    x, y = xy
    for line in text_wrap(draw, text, font, max_width):
        draw.text((x, y), line, font=font, fill=fill)
        y += draw.textbbox((0, 0), line, font=font)[3] + line_gap
    return y


def rounded(draw: ImageDraw.ImageDraw, box, radius: int, fill, outline=None, width=1):
    draw.rounded_rectangle(box, radius=radius, fill=fill, outline=outline, width=width)


def base_canvas() -> Image.Image:
    img = Image.new("RGB", (W, H), PAPER)
    draw = ImageDraw.Draw(img)
    for y in range(H):
        t = y / H
        r = int(251 * (1 - t) + 246 * t)
        g = int(252 * (1 - t) + 248 * t)
        b = int(254 * (1 - t) + 252 * t)
        draw.line([(0, y), (W, y)], fill=(r, g, b))
    return img


def fit_image(img: Image.Image, box: tuple[int, int, int, int], cover: bool = False) -> Image.Image:
    bw, bh = box[2] - box[0], box[3] - box[1]
    iw, ih = img.size
    scale = max(bw / iw, bh / ih) if cover else min(bw / iw, bh / ih)
    resized = img.resize((int(iw * scale), int(ih * scale)), Image.LANCZOS)
    if cover:
        left = max(0, (resized.width - bw) // 2)
        top = max(0, (resized.height - bh) // 2)
        resized = resized.crop((left, top, left + bw, top + bh))
    return resized


def paste_center(canvas: Image.Image, img: Image.Image, box, cover: bool = False):
    fitted = fit_image(img, box, cover=cover)
    x = box[0] + ((box[2] - box[0]) - fitted.width) // 2
    y = box[1] + ((box[3] - box[1]) - fitted.height) // 2
    canvas.paste(fitted, (x, y))


def load_assets() -> dict[str, Image.Image]:
    return {
        "logo": Image.open(ASSETS / "branding" / "vla_corrector_logo.png").convert("RGB"),
        "lab": Image.open(ASSETS / "branding" / "lab_logo.jpg").convert("RGB"),
        "teaser": Image.open(ASSETS / "images" / "teaser_open_loop_vs_closed_loop.png").convert("RGB"),
        "method": Image.open(ASSETS / "images" / "method_overview.png").convert("RGB"),
        "results": Image.open(ASSETS / "images" / "results_pareto.png").convert("RGB"),
        "truncation": Image.open(ASSETS / "images" / "truncation_phase_analysis.png").convert("RGB"),
        "recovery": Image.open(ASSETS / "images" / "qualitative_recovery.png").convert("RGB"),
    }


def extract_video_thumb(video: Path, out: Path, seconds: float = 4.0) -> None:
    out.parent.mkdir(parents=True, exist_ok=True)
    subprocess.run(
        [
            "ffmpeg",
            "-y",
            "-ss",
            str(seconds),
            "-i",
            str(video),
            "-frames:v",
            "1",
            "-q:v",
            "3",
            str(out),
        ],
        check=True,
        stdout=subprocess.DEVNULL,
        stderr=subprocess.DEVNULL,
    )


def prepare_thumbnails() -> list[Path]:
    videos = [
        ASSETS / "videos" / "drawer_alignment_perturbation.mp4",
        ASSETS / "videos" / "block_to_blue_bowl_perturbation.mp4",
        ASSETS / "videos" / "block_to_white_bowl_perturbation.mp4",
    ]
    names = ["drawer_alignment.jpg", "blue_bowl.jpg", "white_bowl.jpg"]
    outs: list[Path] = []
    for video, name in zip(videos, names):
        out = THUMBS / name
        extract_video_thumb(video, out)
        outs.append(out)
    return outs


def header(draw: ImageDraw.ImageDraw, section: str, title: str):
    draw.text((80, 52), section.upper(), font=pil_font(22, bold=True), fill=BLUE)
    draw.text((80, 88), title, font=pil_font(50, bold=True, serif=True), fill=TEXT)


def chip(draw: ImageDraw.ImageDraw, xy, label: str):
    x, y = xy
    f = pil_font(19, bold=True)
    bbox = draw.textbbox((0, 0), label, font=f)
    width = bbox[2] - bbox[0] + 34
    rounded(draw, (x, y, x + width, y + 42), 21, BLUE_SOFT)
    draw.text((x + 17, y + 10), label, font=f, fill=BLUE)
    return x + width + 14


def result_card(draw, box, title: str, value: str, desc: str, color=GREEN):
    rounded(draw, box, 28, "#ffffff", LINE, 2)
    x, y = box[0] + 26, box[1] + 26
    draw.text((x, y), title, font=pil_font(24, bold=True), fill=TEXT)
    draw.text((x, y + 54), value, font=pil_font(52, bold=True, serif=True), fill=color)
    draw_wrapped(draw, (x, y + 122), desc, pil_font(22), MUTED, box[2] - box[0] - 52, 6)


def slide_title(assets) -> Image.Image:
    img = base_canvas()
    draw = ImageDraw.Draw(img)
    paste_center(img, assets["logo"], (90, 95, 330, 335))
    draw.text((380, 105), "VLA-Corrector", font=pil_font(88, bold=True, serif=True), fill=BLUE_DARK)
    draw.text(
        (384, 218),
        "Lightweight Detect-and-Correct Inference\nfor Adaptive Action Horizon",
        font=pil_font(40, bold=True),
        fill=TEXT,
        spacing=12,
    )
    x = 384
    for label in ["Embodied AI", "Vision-Language-Action", "Action Correction", "~40M Corrector"]:
        x = chip(draw, (x, 392), label)
    rounded(draw, (110, 565, 1490, 760), 32, "#ffffff", LINE, 2)
    draw.text((150, 610), "Motivation", font=pil_font(28, bold=True), fill=BLUE)
    draw_wrapped(
        draw,
        (150, 660),
        "Can a frozen VLA keep long-horizon efficiency while recovering from execution drift before stale actions accumulate?",
        pil_font(38, bold=True),
        TEXT,
        1230,
        10,
    )
    draw.text(
        (110, 815),
        "Code: github.com/ZJU-OmniAI/vla-corrector   ·   arXiv: arxiv.org/abs/2607.01804",
        font=pil_font(22),
        fill=MUTED,
    )
    return img


def slide_problem(assets) -> Image.Image:
    img = base_canvas()
    draw = ImageDraw.Draw(img)
    header(draw, "01  Motivation", "Open-loop VLA execution leaves blind spots.")
    rounded(draw, (80, 165, 1520, 610), 22, "#ffffff", LINE, 2)
    paste_center(img, assets["teaser"], (110, 190, 1490, 580))
    rounded(draw, (120, 650, 760, 805), 22, "#ffffff", LINE, 2)
    rounded(draw, (840, 650, 1480, 805), 22, "#ffffff", LINE, 2)
    draw.text((150, 682), "Long action chunks", font=pil_font(28, bold=True), fill=TEXT)
    draw_wrapped(draw, (150, 730), "Efficient VLA calls, but queued actions can become stale after drift.", pil_font(25), MUTED, 540)
    draw.text((870, 682), "Per-step replanning", font=pil_font(28, bold=True), fill=TEXT)
    draw_wrapped(draw, (870, 730), "More reactive, but expensive for large VLA policies.", pil_font(25), MUTED, 540)
    return img


def slide_method(assets) -> Image.Image:
    img = base_canvas()
    draw = ImageDraw.Draw(img)
    header(draw, "02  Method", "Detect drift, truncate stale chunks, guide recovery.")
    rounded(draw, (80, 165, 1520, 660), 22, "#ffffff", LINE, 2)
    paste_center(img, assets["method"], (110, 195, 1490, 628))
    bullets = [
        "Frozen VLA backbone with an external latent dynamics corrector.",
        "Latent-space Vision Monitor detects persistent visual dynamics mismatch.",
        "Online Gradient Guidance is applied only to the recovery query.",
    ]
    x = 100
    for i, text in enumerate(bullets, 1):
        rounded(draw, (x, 700, x + 445, 825), 22, "#ffffff", LINE, 2)
        draw.text((x + 24, 728), f"{i:02d}", font=pil_font(24, bold=True), fill=BLUE)
        draw_wrapped(draw, (x + 82, 724), text, pil_font(23, bold=True), TEXT, 325, 5)
        x += 490
    return img


def slide_corrector(assets) -> Image.Image:
    img = base_canvas()
    draw = ImageDraw.Draw(img)
    header(draw, "03  Corrector", "A small external model instead of full VLA retraining.")
    rounded(draw, (80, 165, 980, 805), 22, "#ffffff", LINE, 2)
    paste_center(img, assets["truncation"], (110, 195, 950, 775))
    result_card(draw, (1040, 180, 1485, 350), "Trainable component", "~40M", "Residual MLP corrector reported at roughly 38--42M parameters.")
    result_card(draw, (1040, 390, 1485, 560), "Policy backbone", "Frozen", "The base VLA is not fully retrained for correction.")
    result_card(draw, (1040, 600, 1485, 770), "Action horizon", "Adaptive", "Stable chunks continue; stale chunks are interrupted.")
    return img


def slide_results(assets) -> Image.Image:
    img = base_canvas()
    draw = ImageDraw.Draw(img)
    header(draw, "04  Results", "Reported gains across simulation and real robots.")
    result_card(draw, (90, 175, 440, 390), "MetaWorld / PI0.5", "+15.65 pts", "48.70% -> 64.35% avg. success")
    result_card(draw, (480, 175, 830, 390), "LIBERO / PI0.5", "+3.80 pts", "94.00% -> 97.80% few-shot success")
    result_card(draw, (870, 175, 1220, 390), "AgileX PiPER", "+17.7 pts", "55.6% -> 73.3% real-world success")
    result_card(draw, (1260, 175, 1535, 390), "Critical phases", "83.7%", "Truncations in manually labeled critical phases", BLUE)
    rounded(draw, (90, 435, 1535, 825), 26, "#ffffff", LINE, 2)
    paste_center(img, assets["results"], (120, 465, 730, 795))
    paste_center(img, assets["recovery"], (780, 485, 1500, 690))
    draw_wrapped(
        draw,
        (810, 724),
        "Qualitative recovery: truncate stale chunk, replan with OGG, complete the task.",
        pil_font(25),
        MUTED,
        630,
        8,
    )
    return img


def slide_demos(assets, thumbs: list[Path]) -> Image.Image:
    img = base_canvas()
    draw = ImageDraw.Draw(img)
    header(draw, "05  Real Robots", "Perturbation demos are kept as silent project-page clips.")
    labels = [
        ("Drawer alignment", "Human moves the drawer during execution."),
        ("Block to blue bowl", "Target bowl is shifted mid-task."),
        ("Block to white bowl", "Object/target perturbation during execution."),
    ]
    x = 85
    for thumb, (title, desc) in zip(thumbs, labels):
        rounded(draw, (x, 190, x + 460, 690), 26, "#ffffff", LINE, 2)
        frame = Image.open(thumb).convert("RGB")
        paste_center(img, frame, (x + 22, 220, x + 438, 455), cover=True)
        draw.text((x + 28, 500), title, font=pil_font(30, bold=True), fill=TEXT)
        draw_wrapped(draw, (x + 28, 552), desc, pil_font(24), MUTED, 390, 8)
        x += 505
    rounded(draw, (145, 735, 1455, 815), 24, BLUE_SOFT, None)
    draw.text((180, 760), "The presentation uses static thumbnails; the project page still hosts the three silent demo videos.", font=pil_font(28, bold=True), fill=BLUE_DARK)
    return img


def slide_resources(assets) -> Image.Image:
    img = base_canvas()
    draw = ImageDraw.Draw(img)
    paste_center(img, assets["logo"], (95, 85, 300, 290))
    draw.text((345, 105), "VLA-Corrector", font=pil_font(72, bold=True, serif=True), fill=BLUE_DARK)
    draw.text((350, 202), "Resources", font=pil_font(40, bold=True), fill=TEXT)
    rows = [
        ("Code", "https://github.com/ZJU-OmniAI/vla-corrector"),
        ("Project page", "https://zju-omniai.github.io/vla-corrector/"),
        ("Paper", "https://arxiv.org/abs/2607.01804"),
        ("arXiv", "https://arxiv.org/abs/2607.01804"),
        ("Editable slides", "docs/assets/presentation/vla_corrector_presentation.pptx"),
    ]
    y = 360
    for key, val in rows:
        rounded(draw, (180, y, 1420, y + 72), 20, "#ffffff", LINE, 2)
        draw.text((215, y + 20), key, font=pil_font(26, bold=True), fill=BLUE)
        draw.text((455, y + 20), val, font=pil_font(26), fill=TEXT)
        y += 92
    return img


def save_slide_images() -> list[Path]:
    PRESENTATION.mkdir(parents=True, exist_ok=True)
    SLIDES.mkdir(parents=True, exist_ok=True)
    assets = load_assets()
    thumbs = prepare_thumbnails()
    images = [
        slide_title(assets),
        slide_problem(assets),
        slide_method(assets),
        slide_corrector(assets),
        slide_results(assets),
        slide_demos(assets, thumbs),
        slide_resources(assets),
    ]
    paths: list[Path] = []
    for idx, image in enumerate(images, 1):
        out = SLIDES / f"slide_{idx:02d}.png"
        image.save(out, optimize=True)
        paths.append(out)
    images[0].save(
        PRESENTATION / "vla_corrector_presentation.pdf",
        save_all=True,
        append_images=images[1:],
        resolution=144,
    )
    return paths


def add_textbox(slide, x, y, w, h, text, size=24, bold=False, color=TEXT, font="Arial", align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    p = frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)
    if align is not None:
        p.alignment = align
    return box


def add_picture_fit(slide, image_path: Path, x, y, w, h):
    img = Image.open(image_path)
    iw, ih = img.size
    box_ratio = w / h
    img_ratio = iw / ih
    if img_ratio > box_ratio:
        width = w
        height = w / img_ratio
        left = x
        top = y + (h - height) / 2
    else:
        height = h
        width = h * img_ratio
        left = x + (w - width) / 2
        top = y
    slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), Inches(width), Inches(height))


def add_chip(slide, x, y, text):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(1.65), Inches(0.36))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(BLUE_SOFT)
    shape.line.fill.background()
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = text
    r.font.name = "Arial"
    r.font.bold = True
    r.font.size = Pt(10.5)
    r.font.color.rgb = rgb(BLUE)


def add_result_box(slide, x, y, w, title, value, desc, color=GREEN):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(1.35))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb("#ffffff")
    shape.line.color.rgb = rgb(LINE)
    add_textbox(slide, x + 0.18, y + 0.16, w - 0.36, 0.25, title, 11, True)
    add_textbox(slide, x + 0.18, y + 0.45, w - 0.36, 0.45, value, 26, True, color, "Georgia")
    add_textbox(slide, x + 0.18, y + 0.98, w - 0.36, 0.3, desc, 9.5, False, MUTED)


def build_pptx() -> Path:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    asset_paths = {
        "logo": ASSETS / "branding" / "vla_corrector_logo.png",
        "teaser": ASSETS / "images" / "teaser_open_loop_vs_closed_loop.png",
        "method": ASSETS / "images" / "method_overview.png",
        "results": ASSETS / "images" / "results_pareto.png",
        "truncation": ASSETS / "images" / "truncation_phase_analysis.png",
        "recovery": ASSETS / "images" / "qualitative_recovery.png",
    }
    thumbs = prepare_thumbnails()

    def background(slide):
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = rgb(PAPER)

    slide = prs.slides.add_slide(blank)
    background(slide)
    slide.shapes.add_picture(str(asset_paths["logo"]), Inches(0.75), Inches(0.7), Inches(2.0), Inches(2.0))
    add_textbox(slide, 3.15, 0.82, 8.7, 0.75, "VLA-Corrector", 50, True, BLUE_DARK, "Georgia")
    add_textbox(slide, 3.2, 1.8, 8.5, 0.86, "Lightweight Detect-and-Correct Inference\nfor Adaptive Action Horizon", 24, True)
    for x, text in [(3.2, "Embodied AI"), (4.9, "VLA"), (6.2, "Action Correction"), (8.15, "~40M Corrector")]:
        add_chip(slide, x, 3.15, text)
    add_textbox(slide, 1.05, 4.85, 11.2, 0.35, "Motivation", 16, True, BLUE)
    add_textbox(
        slide,
        1.05,
        5.35,
        11.3,
        0.85,
        "Can a frozen VLA keep long-horizon efficiency while recovering from execution drift before stale actions accumulate?",
        25,
        True,
    )
    add_textbox(
        slide,
        0.9,
        6.85,
        11.7,
        0.25,
        "Code: github.com/ZJU-OmniAI/vla-corrector   ·   arXiv: arxiv.org/abs/2607.01804",
        11,
        False,
        MUTED,
    )

    slide = prs.slides.add_slide(blank)
    background(slide)
    add_textbox(slide, 0.65, 0.42, 2.2, 0.24, "01  MOTIVATION", 12, True, BLUE)
    add_textbox(slide, 0.65, 0.72, 11.6, 0.5, "Open-loop VLA execution leaves blind spots.", 30, True, TEXT, "Georgia")
    add_picture_fit(slide, asset_paths["teaser"], 0.9, 1.42, 11.55, 3.4)
    add_textbox(slide, 1.05, 5.45, 4.9, 0.28, "Long action chunks", 17, True)
    add_textbox(slide, 1.05, 5.85, 4.9, 0.58, "Efficient VLA calls, but queued actions can become stale after drift.", 15, False, MUTED)
    add_textbox(slide, 7.1, 5.45, 4.9, 0.28, "Per-step replanning", 17, True)
    add_textbox(slide, 7.1, 5.85, 4.9, 0.58, "More reactive, but expensive for large VLA policies.", 15, False, MUTED)

    slide = prs.slides.add_slide(blank)
    background(slide)
    add_textbox(slide, 0.65, 0.42, 1.8, 0.24, "02  METHOD", 12, True, BLUE)
    add_textbox(slide, 0.65, 0.72, 11.8, 0.5, "Detect drift, truncate stale chunks, guide recovery.", 30, True, TEXT, "Georgia")
    add_picture_fit(slide, asset_paths["method"], 0.8, 1.45, 11.85, 4.0)
    add_textbox(slide, 0.95, 5.86, 3.55, 0.64, "Frozen VLA backbone with an external latent dynamics corrector.", 13.5, True)
    add_textbox(slide, 4.85, 5.86, 3.65, 0.64, "Latent-space Vision Monitor detects persistent visual dynamics mismatch.", 13.5, True)
    add_textbox(slide, 8.75, 5.86, 3.65, 0.64, "Online Gradient Guidance is applied only to the recovery query.", 13.5, True)

    slide = prs.slides.add_slide(blank)
    background(slide)
    add_textbox(slide, 0.65, 0.42, 2.1, 0.24, "03  CORRECTOR", 12, True, BLUE)
    add_textbox(slide, 0.65, 0.72, 11.7, 0.5, "A small external model instead of full VLA retraining.", 30, True, TEXT, "Georgia")
    add_picture_fit(slide, asset_paths["truncation"], 0.75, 1.45, 7.2, 5.45)
    add_result_box(slide, 8.6, 1.65, 3.55, "Trainable component", "~40M", "Residual MLP corrector at roughly 38--42M parameters.")
    add_result_box(slide, 8.6, 3.35, 3.55, "Policy backbone", "Frozen", "The base VLA is not fully retrained.")
    add_result_box(slide, 8.6, 5.05, 3.55, "Action horizon", "Adaptive", "Stable chunks continue; stale chunks are interrupted.")

    slide = prs.slides.add_slide(blank)
    background(slide)
    add_textbox(slide, 0.65, 0.42, 1.8, 0.24, "04  RESULTS", 12, True, BLUE)
    add_textbox(slide, 0.65, 0.72, 11.7, 0.5, "Reported gains across simulation and real robots.", 30, True, TEXT, "Georgia")
    add_result_box(slide, 0.75, 1.4, 2.75, "MetaWorld / PI0.5", "+15.65 pts", "48.70% -> 64.35% avg. success")
    add_result_box(slide, 3.75, 1.4, 2.75, "LIBERO / PI0.5", "+3.80 pts", "94.00% -> 97.80% few-shot success")
    add_result_box(slide, 6.75, 1.4, 2.75, "AgileX PiPER", "+17.7 pts", "55.6% -> 73.3% real-world success")
    add_result_box(slide, 9.75, 1.4, 2.55, "Critical phases", "83.7%", "Truncations in manually labeled critical phases", BLUE)
    add_picture_fit(slide, asset_paths["results"], 1.0, 3.55, 5.2, 2.7)
    add_picture_fit(slide, asset_paths["recovery"], 6.45, 3.75, 5.75, 1.85)
    add_textbox(slide, 6.7, 5.95, 5.2, 0.38, "Qualitative recovery: truncate stale chunk, replan with OGG, complete the task.", 14, False, MUTED)

    slide = prs.slides.add_slide(blank)
    background(slide)
    add_textbox(slide, 0.65, 0.42, 2.3, 0.24, "05  REAL ROBOTS", 12, True, BLUE)
    add_textbox(slide, 0.65, 0.72, 11.7, 0.5, "Perturbation demos are kept as silent project-page clips.", 30, True, TEXT, "Georgia")
    for x, thumb, title, desc in [
        (0.7, thumbs[0], "Drawer alignment", "Human moves the drawer during execution."),
        (4.65, thumbs[1], "Block to blue bowl", "Target bowl is shifted mid-task."),
        (8.6, thumbs[2], "Block to white bowl", "Object/target perturbation during execution."),
    ]:
        add_picture_fit(slide, thumb, x, 1.75, 3.55, 2.25)
        add_textbox(slide, x, 4.25, 3.55, 0.35, title, 18, True)
        add_textbox(slide, x, 4.72, 3.55, 0.72, desc, 14, False, MUTED)
    add_textbox(slide, 1.2, 6.45, 10.9, 0.35, "Static thumbnails in the slide deck; full clips remain silent on the project page.", 17, True, BLUE_DARK, align=PP_ALIGN.CENTER)

    slide = prs.slides.add_slide(blank)
    background(slide)
    slide.shapes.add_picture(str(asset_paths["logo"]), Inches(0.75), Inches(0.65), Inches(1.75), Inches(1.75))
    add_textbox(slide, 2.85, 0.85, 8.6, 0.66, "VLA-Corrector", 42, True, BLUE_DARK, "Georgia")
    add_textbox(slide, 2.9, 1.68, 4.0, 0.35, "Resources", 24, True)
    for y, key, val in [
        (3.05, "Code", "https://github.com/ZJU-OmniAI/vla-corrector"),
        (3.82, "Project page", "https://zju-omniai.github.io/vla-corrector/"),
        (4.59, "Paper", "https://arxiv.org/abs/2607.01804"),
        (5.36, "arXiv", "https://arxiv.org/abs/2607.01804"),
        (6.13, "Editable slides", "docs/assets/presentation/vla_corrector_presentation.pptx"),
    ]:
        add_textbox(slide, 1.65, y, 2.2, 0.25, key, 16, True, BLUE)
        add_textbox(slide, 4.05, y, 7.9, 0.25, val, 16, False)

    out = PRESENTATION / "vla_corrector_presentation.pptx"
    prs.save(out)
    return out


def main() -> None:
    PRESENTATION.mkdir(parents=True, exist_ok=True)
    slide_paths = save_slide_images()
    pptx = build_pptx()
    print(f"Wrote {pptx}")
    print(f"Wrote {PRESENTATION / 'vla_corrector_presentation.pdf'}")
    for slide in slide_paths:
        print(slide)


if __name__ == "__main__":
    main()
